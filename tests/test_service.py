"""Unit tests for extract.service."""

import io
import zipfile
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

from extract.service import (
    DEFAULT_IMAGE_EXTENSIONS,
    DEFAULT_VIDEO_EXTENSIONS,
    MediaInfo,
    PowerPointMediaExtractor,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_mock_presentation(shapes_per_slide=None):
    """Return a mock Presentation object.

    shapes_per_slide: list of lists of shape mocks, one inner list per slide.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = MagicMock()
    slides = []
    for shapes in (shapes_per_slide or []):
        slide = MagicMock()
        slide.shapes = shapes
        slides.append(slide)
    prs.slides = slides
    return prs


def _make_picture_shape(shape_id: int):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    shape = MagicMock()
    shape.shape_id = shape_id
    shape.shape_type = MSO_SHAPE_TYPE.PICTURE
    return shape


def _make_video_shape(shape_id: int):
    from pptx.enum.shapes import PP_MEDIA_TYPE

    shape = MagicMock()
    shape.shape_id = shape_id
    shape.media_type = PP_MEDIA_TYPE.MOVIE
    shape.shape_type = None
    return shape


def _make_pptx_zip(media_files: list[str]) -> bytes:
    """Return bytes of a minimal zip that mimics a .pptx archive."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        for name in media_files:
            zf.writestr(name, b"fake-media-content")
    return buf.getvalue()


# ---------------------------------------------------------------------------
# MediaInfo
# ---------------------------------------------------------------------------

class TestMediaInfo:
    def test_fields(self):
        info = MediaInfo(shape_id=1, filename="image1", slide_number=2)
        assert info.shape_id == 1
        assert info.filename == "image1"
        assert info.slide_number == 2


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

class TestConstants:
    def test_default_image_extensions_contains_png(self):
        assert "png" in DEFAULT_IMAGE_EXTENSIONS

    def test_default_video_extensions_contains_mp4(self):
        assert "mp4" in DEFAULT_VIDEO_EXTENSIONS


# ---------------------------------------------------------------------------
# PowerPointMediaExtractor.__init__
# ---------------------------------------------------------------------------

class TestInit:
    @patch("extract.service.Presentation")
    def test_invalid_media_type_raises(self, mock_prs, tmp_path):
        mock_prs.return_value = MagicMock()
        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(_make_pptx_zip([]))
        with pytest.raises(ValueError, match="Invalid media_type"):
            PowerPointMediaExtractor(pptx, media_type="audio")

    @patch("extract.service.Presentation")
    def test_default_extensions_image(self, mock_prs, tmp_path):
        mock_prs.return_value = MagicMock()
        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(_make_pptx_zip([]))
        extractor = PowerPointMediaExtractor(pptx, media_type="image")
        assert extractor._extensions == list(DEFAULT_IMAGE_EXTENSIONS)

    @patch("extract.service.Presentation")
    def test_default_extensions_video(self, mock_prs, tmp_path):
        mock_prs.return_value = MagicMock()
        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(_make_pptx_zip([]))
        extractor = PowerPointMediaExtractor(pptx, media_type="video")
        assert extractor._extensions == list(DEFAULT_VIDEO_EXTENSIONS)

    @patch("extract.service.Presentation")
    def test_custom_extensions(self, mock_prs, tmp_path):
        mock_prs.return_value = MagicMock()
        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(_make_pptx_zip([]))
        extractor = PowerPointMediaExtractor(pptx, extensions=["GIF", "WEBP"])
        assert extractor._extensions == ["gif", "webp"]

    @patch("extract.service.Presentation")
    def test_output_dir_defaults_to_temp(self, mock_prs, tmp_path):
        mock_prs.return_value = MagicMock()
        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(_make_pptx_zip([]))
        extractor = PowerPointMediaExtractor(pptx)
        assert extractor.output_dir == Path("temp")


# ---------------------------------------------------------------------------
# extract_all_media
# ---------------------------------------------------------------------------

class TestExtractAllMedia:
    @patch("extract.service.Presentation")
    def test_extracts_media_files(self, mock_prs, tmp_path):
        mock_prs.return_value = MagicMock()
        media_names = ["ppt/media/image1.png", "ppt/media/image2.jpg"]
        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(_make_pptx_zip(media_names))

        out = tmp_path / "out"
        extractor = PowerPointMediaExtractor(pptx, output_dir=out)
        count = extractor.extract_all_media()

        assert count == 2

    @patch("extract.service.Presentation")
    def test_skips_non_media_entries(self, mock_prs, tmp_path):
        mock_prs.return_value = MagicMock()
        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(_make_pptx_zip(["ppt/slides/slide1.xml", "docProps/app.xml"]))

        out = tmp_path / "out"
        extractor = PowerPointMediaExtractor(pptx, output_dir=out)
        count = extractor.extract_all_media()

        assert count == 0

    @patch("extract.service.Presentation")
    def test_returns_zero_when_no_media(self, mock_prs, tmp_path):
        mock_prs.return_value = MagicMock()
        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(_make_pptx_zip([]))

        extractor = PowerPointMediaExtractor(pptx, output_dir=tmp_path / "out")
        assert extractor.extract_all_media() == 0


# ---------------------------------------------------------------------------
# extract_filtered_media
# ---------------------------------------------------------------------------

class TestExtractFilteredMedia:
    @patch("extract.service.Presentation")
    def test_extracts_images_by_slide(self, mock_prs, tmp_path):
        picture = _make_picture_shape(shape_id=1)
        prs = _make_mock_presentation([[picture]])
        mock_prs.return_value = prs

        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(_make_pptx_zip(["ppt/media/image1.png"]))

        out = tmp_path / "out"
        extractor = PowerPointMediaExtractor(pptx, media_type="image", output_dir=out)
        count = extractor.extract_filtered_media()

        assert count == 1

    @patch("extract.service.Presentation")
    def test_skips_unmatched_extension(self, mock_prs, tmp_path):
        picture = _make_picture_shape(shape_id=1)
        prs = _make_mock_presentation([[picture]])
        mock_prs.return_value = prs

        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(_make_pptx_zip(["ppt/media/image1.gif"]))

        out = tmp_path / "out"
        extractor = PowerPointMediaExtractor(pptx, media_type="image", output_dir=out)
        count = extractor.extract_filtered_media()

        assert count == 0

    @patch("extract.service.Presentation")
    def test_returns_zero_when_no_shapes(self, mock_prs, tmp_path):
        prs = _make_mock_presentation([[]])
        mock_prs.return_value = prs

        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(_make_pptx_zip(["ppt/media/image1.png"]))

        extractor = PowerPointMediaExtractor(pptx, media_type="image", output_dir=tmp_path / "out")
        assert extractor.extract_filtered_media() == 0
