import logging
from dataclasses import dataclass
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_MEDIA_TYPE

logger = logging.getLogger(__name__)

DEFAULT_PPT_MEDIA_PATH = "ppt/media"
DEFAULT_IMAGE_EXTENSIONS = ["png", "jpeg", "jpg", "bmp", "svg"]
DEFAULT_VIDEO_EXTENSIONS = ["mp4", "avi", "mpg", "mpeg", "wmv"]


@dataclass
class MediaInfo:
    shape_id: int
    filename: str
    slide_number: int


class PowerPointMediaExtractor:
    def __init__(
        self,
        filepath: str | Path,
        media_type: str = "image",
        output_dir: str | Path = "temp",
        extensions: list[str] | None = None,
    ) -> None:
        self.filepath = Path(filepath)
        self.output_dir = Path(output_dir)
        self.media_type = media_type.lower()

        if self.media_type not in ("image", "video"):
            raise ValueError("Invalid media_type. Use 'image' or 'video'.")

        self.presentation = Presentation(str(self.filepath))

        if extensions:
            self._extensions = [e.lower() for e in extensions]
        elif self.media_type == "image":
            self._extensions = list(DEFAULT_IMAGE_EXTENSIONS)
        else:
            self._extensions = list(DEFAULT_VIDEO_EXTENSIONS)

        self._infos: list[MediaInfo] = []

    def _collect_media_info(self) -> None:
        self._infos = []
        counter = 0
        prefix = "image" if self.media_type == "image" else "media"

        for slide_num, slide in enumerate(self.presentation.slides, start=1):
            for shape in slide.shapes:
                is_match = (
                    getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
                    if self.media_type == "image"
                    else getattr(shape, "media_type", None) == PP_MEDIA_TYPE.MOVIE
                )
                if is_match:
                    counter += 1
                    self._infos.append(
                        MediaInfo(
                            shape_id=shape.shape_id,
                            filename=f"{prefix}{counter}",
                            slide_number=slide_num,
                        )
                    )

    def _find_slide_for_filename(self, filename_stem: str) -> int | None:
        for item in self._infos:
            if item.filename == filename_stem:
                return item.slide_number
        return None

    def extract_all_media(self) -> int:
        """Extract all media from the presentation without filtering."""
        extracted = 0
        with ZipFile(self.filepath, "r") as archive:
            for name in archive.namelist():
                if name.startswith(DEFAULT_PPT_MEDIA_PATH):
                    self.output_dir.mkdir(parents=True, exist_ok=True)
                    archive.extract(name, self.output_dir)
                    extracted += 1

        status = "Completed" if extracted else "Not Found!"
        logger.info("%s — %d media extracted", status, extracted)
        return extracted

    def extract_filtered_media(self) -> int:
        """Extract media filtered by type and extension, organized by slide number."""
        self._collect_media_info()
        if not self._infos:
            logger.info("No media found.")
            return 0

        extracted = 0
        with ZipFile(self.filepath, "r") as archive:
            for name in archive.namelist():
                if not name.startswith(DEFAULT_PPT_MEDIA_PATH):
                    continue
                stem, ext = Path(name).stem, Path(name).suffix
                if ext.lstrip(".").lower() not in self._extensions:
                    continue
                slide_number = self._find_slide_for_filename(stem)
                if slide_number is not None:
                    target_path = self.output_dir / str(slide_number)
                    target_path.mkdir(parents=True, exist_ok=True)
                    archive.extract(name, target_path)
                    extracted += 1

        status = "Completed" if extracted else "Not Found!"
        logger.info("%s — %d media extracted", status, extracted)
        return extracted
